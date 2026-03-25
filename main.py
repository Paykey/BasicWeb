import os
import httpx
import io
import tempfile
from pathlib import Path
from fastapi import FastAPI, File, UploadFile, HTTPException
from pydantic import BaseModel
from dotenv import load_dotenv
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader
from pptx import Presentation

# from openai import OpenAI

load_dotenv() # .env 파일에 있는 환경변수 불러오기

# 프로젝트 루트 경로 기준으로 static/templates 파일을 안전하게 참조
BASE_DIR = Path(__file__).resolve().parent

app = FastAPI()
# 정적 파일(CSS/JS) 서빙 경로
app.mount("/static", StaticFiles(directory=BASE_DIR / "static"), name="static")

# OpenAI API 사용 시 아래 두 줄 주석 해제
# from openai import OpenAI
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

LOCAL_LLM_ENDPOINT = os.getenv("LOCAL_LLM_ENDPOINT", "http://localhost:11434/api/generate")
LOCAL_LLM_MODEL = os.getenv("LOCAL_LLM_MODEL", "qwen2.5:14b")

@app.get("/")
def root():
    # 첫 화면으로 HTML 템플릿 반환
    return FileResponse(BASE_DIR / "templates" / "index.html")

# PDF/PPTX 파일에서 텍스트 추출 함수
def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages_text = []
    for page in reader.pages:
        pages_text.append(page.extract_text() or "")
    return "\n".join(pages_text).strip()


def extract_text_from_pptx(temp_path: Path) -> str:
    presentation = Presentation(str(temp_path))
    slide_texts = []
    for slide in presentation.slides:
        shape_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                shape_texts.append(shape.text)
        if shape_texts:
            slide_texts.append("\n".join(shape_texts))
    return "\n\n".join(slide_texts).strip()


# 요약 요청을 처리하는 API 엔드포인트
# 프론트엔드에서 긴 텍스트를 보내면 로컬 LLM으로 요약 결과 반환
@app.post("/summarize")
async def summarize_file(file: UploadFile = File(...)):
    allowed_types = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.ms-powerpoint": "ppt",
    }
    extension = Path(file.filename or "").suffix.lower()
    detected_type = allowed_types.get(file.content_type)

    if extension not in {".pdf", ".pptx", ".ppt"} and detected_type not in {"pdf", "pptx", "ppt"}:
        raise HTTPException(status_code=400, detail="PDF 또는 PPT(PPTX) 파일만 업로드할 수 있습니다.")

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="업로드된 파일이 비어 있습니다.")

    temp_file_path: Path | None = None
    try:
        if extension == ".pdf" or detected_type == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif extension == ".pptx" or detected_type == "pptx":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=BASE_DIR) as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = Path(temp_file.name)
            text = extract_text_from_pptx(temp_file_path)
        else:
            raise HTTPException(
                status_code=400,
                detail="기존 PPT(.ppt) 형식은 현재 지원하지 않습니다. PPTX로 변환 후 업로드해 주세요.",
            )

        if not text:
            raise HTTPException(status_code=400, detail="파일에서 추출된 텍스트가 없습니다.")
    finally:
        if temp_file_path and temp_file_path.exists():
            temp_file_path.unlink()
            
    # 모델에 전달할 역할/출력 형식 지시문 + 사용자 원문 결합
    prompt = (
        "너는 자료를 핵심만 요약해 주는 역할이야. 각 줄은 간결하게 작성해."
        "반드시 한국어로 대답해.\n\n"
        f"[강의자료]\n{text}"
    )

    # 로컬 LLM 호출
    try:
        # Ollama generate API 형식으로 요청
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                LOCAL_LLM_ENDPOINT,
                json={
                    "model": LOCAL_LLM_MODEL,
                    "prompt": prompt,
                    "stream": False,
                },
            )
        response.raise_for_status()
        data = response.json()
        # Ollama 응답의 실제 생성 텍스트는 response 키에 담김
        summary = data.get("response", "요약 결과가 비어 있습니다.")
        return {"summary": summary}
    except Exception as exc:
        # 프론트에서 바로 처리할 수 있도록 HTTP 에러로 변환
        raise HTTPException(status_code=500, detail=f"로컬 LLM 호출 실패: {exc}")

    # OpenAI API로 다시 전환할 때 아래 블록 사용
    # response = client.chat.completions.create(
    #     model="gpt-4o-mini",
    #     messages=[
    #         {
    #             "role": "system",
    #             "content": "너는 대학생의 전공 강의자료를 핵심만 3줄로 요약해 주는 역할이야. 반드시 한국어로 대답해.",
    #         },
    #         {"role": "user", "content": request.text},
    #     ],
    # )
    # return {"summary": response.choices[0].message.content}